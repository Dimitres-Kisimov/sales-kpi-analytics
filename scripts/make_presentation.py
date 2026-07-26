"""make_presentation.py — build the company-facing Executive Business Review.

This is the deliverable I'd actually hand to a leadership team: a clean,
one-chart-per-slide review built straight from `deliverables/analysis.json`.

    python scripts/make_presentation.py

Always writes:
    deliverables/executive_review.pdf     (matplotlib PdfPages, ~8 slides)

Also writes, only if python-pptx is installed:
    deliverables/executive_review.pptx    (same story, charts embedded as PNGs)

If python-pptx isn't there it just skips the .pptx and says so — the PDF is the
guaranteed output.

Author: Dimitres Kisimov.
"""
from __future__ import annotations

import json
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages

ROOT = Path(__file__).resolve().parents[1]
ANALYSIS = ROOT / "deliverables" / "analysis.json"
OUT_PDF = ROOT / "deliverables" / "executive_review.pdf"
OUT_PPTX = ROOT / "deliverables" / "executive_review.pptx"
PNG_DIR = ROOT / "deliverables" / "_slides"

# palette
BLUE = "#2f6bff"
PINK = "#ea4b71"
INK = "#1a1f2b"
MUTE = "#8b8f99"
GRID = "#e6e8ee"
AQUA = "#1baf7a"
GOLD = "#eda100"
PANEL = "#f5f7fb"

plt.rcParams.update({
    "font.family": "DejaVu Sans",
    "figure.facecolor": "white",
    "axes.facecolor": "white",
    "axes.edgecolor": MUTE,
    "text.color": INK,
    "axes.labelcolor": INK,
    "xtick.color": MUTE,
    "ytick.color": MUTE,
})

PAGE = (13.33, 7.5)   # 16:9 widescreen inches


def eur(v: float, mm: bool = False) -> str:
    if mm:
        return f"€{v / 1e6:.2f}M"
    return "€" + f"{round(v):,.0f}"


def new_slide(title: str, subtitle: str = ""):
    fig = plt.figure(figsize=PAGE, dpi=150)
    fig.patch.set_facecolor("white")
    # top rule + title
    fig.text(0.06, 0.90, title, fontsize=26, fontweight="bold", color=INK)
    if subtitle:
        fig.text(0.06, 0.855, subtitle, fontsize=13, color=MUTE)
    fig.add_artist(plt.Line2D([0.06, 0.94], [0.83, 0.83], color=BLUE, lw=2.5))
    return fig


def footer(fig, page: int, total: int):
    fig.text(0.06, 0.04, "Sales & Demand Analytics — Executive Review", fontsize=9, color=MUTE)
    fig.text(0.94, 0.04, f"{page} / {total}", fontsize=9, color=MUTE, ha="right")


def ax_clean(ax):
    for s in ("top", "right"):
        ax.spines[s].set_visible(False)
    ax.spines["left"].set_color(GRID)
    ax.spines["bottom"].set_color(GRID)
    ax.grid(axis="y", color=GRID, lw=0.8)
    ax.set_axisbelow(True)


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #
def slide_title(a):
    fig = plt.figure(figsize=PAGE, dpi=150)
    fig.patch.set_facecolor(INK)
    fig.add_artist(plt.Line2D([0.08, 0.30], [0.60, 0.60], color=BLUE, lw=4))
    fig.text(0.08, 0.66, "Sales & Demand Analytics", fontsize=40, fontweight="bold", color="white")
    fig.text(0.08, 0.52, "Executive Business Review", fontsize=24, color="#c8ccd6")
    k = a["kpi"]
    g = a["growth"]
    yoy = g.get("yoy_pct")
    strip = f"{eur(k['revenue_eur'], mm=True)} revenue   ·   {k['gross_margin_pct']:.1f}% margin"
    if yoy is not None:
        strip += f"   ·   {yoy:+.1f}% YoY"
    strip += f"   ·   OTIF {k['otif_pct']:.0f}%"
    fig.text(0.08, 0.40, strip, fontsize=15, color=PINK, fontweight="bold")
    fig.text(0.08, 0.12, "Prepared by Dimitres Kisimov", fontsize=13, color="#c8ccd6")
    fig.text(0.08, 0.08, "Synthetic distributor dataset · 24 months", fontsize=11, color=MUTE)
    return fig


def slide_scorecard(a):
    fig = new_slide("KPI scorecard", "Headline commercial and service metrics")
    k, g = a["kpi"], a["growth"]
    yoy = g.get("yoy_pct")
    cards = [
        ("Revenue", eur(k["revenue_eur"], mm=True), f"{k['orders']:,} orders"),
        ("Gross margin", f"{k['gross_margin_pct']:.1f}%", eur(k["gross_margin_eur"])),
        ("Revenue YoY", f"{yoy:+.1f}%" if yoy is not None else "—", "12m vs prior 12m"),
        ("Avg order value", eur(k["aov_eur"]), f"{k['active_customers']} customers"),
        ("OTIF", f"{k['otif_pct']:.0f}%", "on-time & in-full"),
        ("Discount leakage", f"{k['discount_leakage_pct']:.1f}%", "of list value"),
    ]
    cols = 3
    x0, y0, w, h, gx, gy = 0.06, 0.44, 0.28, 0.28, 0.02, 0.06
    for i, (label, value, foot) in enumerate(cards):
        cx = x0 + (i % cols) * (w + gx)
        cy = y0 - (i // cols) * (h + gy)
        ax = fig.add_axes([cx, cy, w, h])
        ax.axis("off")
        ax.add_patch(plt.Rectangle((0, 0), 1, 1, facecolor=PANEL, edgecolor=GRID, lw=1))
        ax.text(0.07, 0.72, label, fontsize=13, color=MUTE, fontweight="bold")
        ax.text(0.07, 0.36, value, fontsize=30, color=INK, fontweight="bold")
        ax.text(0.07, 0.14, foot, fontsize=11, color=MUTE)
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
    return fig


def slide_forecast(a):
    fig = new_slide("Revenue trend & forecast",
                    "Monthly revenue with a 3-month forecast (model chosen by rolling-origin CV)")
    rf = a["revenue_forecast"]
    hist, fc = rf["history"], rf["forecast"]
    ax = fig.add_axes([0.08, 0.14, 0.84, 0.60])
    ax_clean(ax)
    xh = list(range(1, len(hist) + 1))
    ax.plot(xh, [v / 1e6 for v in hist], color=BLUE, lw=2.6, label="actual")
    xf = list(range(len(hist), len(hist) + len(fc) + 1))
    ax.plot(xf, [hist[-1] / 1e6] + [v / 1e6 for v in fc], color=PINK, lw=2.6, ls="--",
            marker="o", ms=7, label=f"forecast ({rf['winner']})")
    ax.set_ylabel("revenue (€M)")
    ax.set_xlabel("month index")
    ax.legend(frameon=False, fontsize=12, loc="upper left")
    ax.set_ylim(bottom=0)
    fig.text(0.08, 0.80, f"Selected model: {rf['winner']}  ·  CV-MASE {rf['cv_mase']:.2f} "
             f"(a value <1 beats the seasonal-naive baseline)", fontsize=12, color=MUTE)
    return fig


def slide_bridge(a):
    fig = new_slide("Margin bridge", "What moved gross margin year-on-year: price / volume / mix")
    b = a["margin_bridge"]
    ax = fig.add_axes([0.10, 0.16, 0.82, 0.56])
    ax_clean(ax)
    if not b.get("available"):
        ax.text(0.5, 0.5, "Not enough history for a YoY bridge", ha="center", color=MUTE)
        ax.axis("off")
        return fig
    start = b["prior_margin_eur"]
    steps = [("Prior\nmargin", start, INK),
             ("Price", b["price_effect_eur"], BLUE),
             ("Volume", b["volume_effect_eur"], AQUA),
             ("Mix", b["mix_effect_eur"], GOLD),
             ("Current\nmargin", None, INK)]
    running = start
    xs = range(len(steps))
    for i, (label, val, color) in enumerate(steps):
        if i == 0:
            ax.bar(i, start / 1e6, color=color, width=0.6)
        elif label.startswith("Current"):
            ax.bar(i, running / 1e6, color=color, width=0.6)
        else:
            bottom = running
            running += val
            lo = min(bottom, running)
            ax.bar(i, abs(val) / 1e6, bottom=lo / 1e6, color=color, width=0.6)
            ax.annotate(f"{val:+,.0f}", (i, max(bottom, running) / 1e6),
                        textcoords="offset points", xytext=(0, 6),
                        ha="center", fontsize=11, color=INK)
    ax.set_xticks(list(xs))
    ax.set_xticklabels([s[0] for s in steps], fontsize=11)
    ax.set_ylabel("gross margin (€M)")
    fig.text(0.10, 0.80, f"Total change {b['total_change_eur']:+,.0f} EUR "
             f"({eur(start, mm=True)} → {eur(b['current_margin_eur'], mm=True)})",
             fontsize=12, color=MUTE)
    return fig


def slide_abcxyz(a):
    fig = new_slide("Portfolio — ABC × XYZ", "Revenue concentration vs demand variability by category")
    ax = fig.add_axes([0.08, 0.14, 0.62, 0.62])
    ax_clean(ax)
    rows = a["abc_xyz"]
    cats = [r["product_category"] for r in rows][::-1]
    revs = [r["revenue_eur"] / 1e6 for r in rows][::-1]
    colmap = {"A": BLUE, "B": AQUA, "C": MUTE}
    colors = [colmap.get(r["abc"], MUTE) for r in rows][::-1]
    ax.barh(cats, revs, color=colors)
    ax.set_xlabel("revenue (€M)")
    for i, r in enumerate(rows[::-1]):
        ax.text(revs[i], i, f"  {r['class']}", va="center", fontsize=10, color=INK)
    # legend panel
    lx = fig.add_axes([0.74, 0.30, 0.20, 0.40])
    lx.axis("off")
    lx.text(0, 1.0, "Class key", fontsize=13, fontweight="bold", color=INK)
    notes = [("A", BLUE, "top ~80% of revenue"), ("B", AQUA, "next ~15%"),
             ("C", MUTE, "long tail"), ("", None, ""),
             ("X", INK, "steady demand"), ("Y", INK, "variable"), ("Z", INK, "erratic")]
    for j, (lab, col, txt) in enumerate(notes):
        yy = 0.86 - j * 0.12
        if col and lab:
            lx.add_patch(plt.Rectangle((0, yy), 0.12, 0.06, color=col))
        if lab:
            lx.text(0.18, yy, f"{lab} — {txt}", fontsize=11, color=INK, va="bottom")
    lx.set_xlim(0, 1)
    lx.set_ylim(0, 1)
    return fig


def slide_spend(a):
    fig = new_slide("Expenditure & spend", "COGS plus the two leaks a P&L usually hides")
    sp = a.get("expenditure")
    ax = fig.add_axes([0.08, 0.16, 0.52, 0.56])
    ax_clean(ax)
    if sp:
        cts = sp["cost_to_serve_by_channel"]
        chans = [c["channel"] for c in cts]
        cogs = [c["cogs_eur"] / 1e6 for c in cts]
        ret = [c["returns_cost_eur"] / 1e6 for c in cts]
        leak = [c["discount_leakage_eur"] / 1e6 for c in cts]
        ax.bar(chans, cogs, color=BLUE, label="COGS")
        ax.bar(chans, ret, bottom=cogs, color=GOLD, label="returns cost")
        ax.bar(chans, leak, bottom=[c + r for c, r in zip(cogs, ret, strict=False)], color=PINK,
               label="discount leakage")
        ax.set_ylabel("€M")
        ax.legend(frameon=False, fontsize=11)
        ax.tick_params(axis="x", labelrotation=15)
        # callout panel
        px = fig.add_axes([0.64, 0.24, 0.30, 0.44])
        px.axis("off")
        px.add_patch(plt.Rectangle((0, 0), 1, 1, facecolor=PANEL, edgecolor=GRID, lw=1))
        px.text(0.08, 0.82, "Where the money goes", fontsize=13, fontweight="bold", color=INK)
        px.text(0.08, 0.62, f"Total COGS  {eur(sp['total_cogs_eur'], mm=True)}", fontsize=13, color=INK)
        px.text(0.08, 0.44, f"Discount leakage  {eur(sp['discount_leakage_eur'])}",
                fontsize=12, color=PINK)
        px.text(0.08, 0.34, f"   = {sp['discount_leakage_pct']:.1f}% of list value",
                fontsize=11, color=MUTE)
        px.text(0.08, 0.16, f"Returns cost  {eur(sp['returns_cost_total_eur'])}",
                fontsize=12, color=GOLD)
        px.set_xlim(0, 1)
        px.set_ylim(0, 1)
    return fig


def slide_leakage(a):
    """The discount-leakage drill-down: waterfall + who/where offender split."""
    fig = new_slide("Discount leakage — who, where",
                    "Gross list value to net revenue; leakage split vs an assumed "
                    "10% sanctioned-discount ceiling")
    dd = (a.get("expenditure") or {}).get("leakage_drilldown")
    if not dd:
        ax = fig.add_axes([0.1, 0.2, 0.8, 0.5])
        ax.axis("off")
        ax.text(0.5, 0.5, "Re-run the analysis to add the drill-down", ha="center", color=MUTE)
        return fig

    # left: waterfall gross -> within-policy -> excess -> net
    ax = fig.add_axes([0.07, 0.14, 0.46, 0.58])
    ax_clean(ax)
    gross = dd["gross_list_value_eur"] / 1e6
    within = dd["within_policy_discount_eur"] / 1e6
    excess = dd["excess_discount_eur"] / 1e6
    net = dd["net_revenue_eur"] / 1e6
    bars = [("Gross\nlist", 0.0, gross, MUTE, f"{gross:,.2f}M"),
            ("Within\npolicy", gross - within, within, BLUE, f"-{within:,.2f}M"),
            ("Excess", net, excess, PINK, f"-{excess:,.2f}M"),
            ("Net\nrevenue", 0.0, net, INK, f"{net:,.2f}M")]
    for i, (_, bottom, height, color, lab) in enumerate(bars):
        ax.bar(i, height, bottom=bottom, color=color, width=0.6,
               edgecolor="white", linewidth=1.5)
        ax.annotate(lab, (i, bottom + height), textcoords="offset points",
                    xytext=(0, 6), ha="center", fontsize=11, color=INK)
    for i, lvl in enumerate([gross, gross - within, net]):
        ax.plot([i + 0.3, i + 0.7], [lvl, lvl], color=MUTE, lw=1, ls=":")
    ax.set_xticks(range(4))
    ax.set_xticklabels([b[0] for b in bars], fontsize=10)
    ax.set_ylabel("€M")

    # right: per-rep offender bars (top 6 by excess)
    reps = dd["by_sales_rep"][:6]
    ax2 = fig.add_axes([0.62, 0.14, 0.32, 0.58])
    ax_clean(ax2)
    names = [f"{r['sales_rep']} ({r['region']})" for r in reps][::-1]
    wp = [r["within_policy_eur"] / 1e3 for r in reps][::-1]
    ex = [r["excess_eur"] / 1e3 for r in reps][::-1]
    ax2.barh(names, wp, color=BLUE, label="within policy", edgecolor="white", linewidth=1)
    ax2.barh(names, ex, left=wp, color=PINK, label="excess", edgecolor="white", linewidth=1)
    for i, r in enumerate(reps[::-1]):
        ax2.text(wp[i] + ex[i], i, f" {r['leakage_eur'] / 1e3:,.0f}k",
                 va="center", fontsize=10, color=INK)
    ax2.set_xlim(0, max(w + e for w, e in zip(wp, ex, strict=False)) * 1.24)
    ax2.set_xlabel("€k leakage")
    ax2.legend(frameon=False, fontsize=9, loc="lower right")
    ax2.tick_params(axis="y", labelsize=10)
    ax2.set_title("Top reps by excess-over-policy", fontsize=12, color=INK,
                  fontweight="bold", loc="left")

    top = dd["by_sales_rep"][0]
    fig.text(0.07, 0.79,
             f"Total leakage {eur(dd['total_leakage_eur'])} vs list · "
             f"{eur(dd['excess_discount_eur'])} above the assumed "
             f"{dd['policy_discount_pct']:.0f}% ceiling · top offender "
             f"{top['sales_rep']} ({top['region']}): {eur(top['leakage_eur'])}, "
             f"{top['orders_above_policy_pct']:.0f}% of orders over policy",
             fontsize=12, color=MUTE)
    return fig


def slide_rfm(a):
    fig = new_slide("Customer segments (RFM)", "Recency / frequency / monetary segmentation")
    segs = a["rfm_segments"]
    order = ["Champions", "Loyal / high-value", "Developing", "At risk", "Churned / dormant"]
    labels = [s for s in order if s in segs]
    vals = [segs[s] for s in labels]
    colmap = {"Champions": BLUE, "Loyal / high-value": AQUA, "Developing": "#6da7ec",
              "At risk": GOLD, "Churned / dormant": PINK}
    colors = [colmap.get(s, MUTE) for s in labels]
    ax = fig.add_axes([0.08, 0.16, 0.54, 0.58])
    ax_clean(ax)
    ax.barh(labels[::-1], vals[::-1], color=colors[::-1])
    ax.set_xlabel("customers")
    at_risk = segs.get("At risk", 0) + segs.get("Churned / dormant", 0)
    total = sum(segs.values())
    px = fig.add_axes([0.66, 0.26, 0.28, 0.40])
    px.axis("off")
    px.add_patch(plt.Rectangle((0, 0), 1, 1, facecolor="#fff0f3", edgecolor=PINK, lw=1.5))
    px.text(0.5, 0.74, "At-risk / churned", fontsize=14, fontweight="bold", color=INK, ha="center")
    px.text(0.5, 0.44, str(at_risk), fontsize=42, fontweight="bold", color=PINK, ha="center")
    px.text(0.5, 0.22, f"of {total} customers — hand to sales\nfor a win-back campaign",
            fontsize=11, color=MUTE, ha="center")
    px.set_xlim(0, 1)
    px.set_ylim(0, 1)
    return fig


def slide_actions(a):
    fig = new_slide("Recommended actions", "The decisions the numbers point to")
    cards = a["decision_cards"]
    ax = fig.add_axes([0.06, 0.10, 0.88, 0.68])
    ax.axis("off")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    n = len(cards)
    top = 0.96
    step = min(0.15, (top - 0.02) / max(n, 1))
    for i, c in enumerate(cards):
        y = top - i * step
        ax.add_patch(plt.Rectangle((0, y - step * 0.75), 0.012, step * 0.7, color=BLUE))
        wrapped = _wrap(c, 96)
        ax.text(0.03, y - step * 0.4, wrapped, fontsize=12.5, color=INK, va="center")
    return fig


def _wrap(text: str, width: int) -> str:
    import textwrap
    return "\n".join(textwrap.wrap(text, width)) or text


SLIDES = [slide_title, slide_scorecard, slide_forecast, slide_bridge,
          slide_abcxyz, slide_spend, slide_leakage, slide_rfm, slide_actions]


def build_pdf(a) -> int:
    total = len(SLIDES)
    with PdfPages(OUT_PDF) as pdf:
        for i, fn in enumerate(SLIDES, start=1):
            fig = fn(a)
            if i > 1:  # no footer on the dark title slide
                footer(fig, i, total)
            pdf.savefig(fig, facecolor=fig.get_facecolor())
            plt.close(fig)
    return total


def build_pptx(a, n_slides: int) -> bool:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return False
    PNG_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for i, fn in enumerate(SLIDES, start=1):
        fig = fn(a)
        png = PNG_DIR / f"slide_{i:02d}.png"
        fig.savefig(png, dpi=150, facecolor=fig.get_facecolor())
        plt.close(fig)
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(OUT_PPTX)
    # tidy up the intermediate PNGs — the .pptx has them embedded now
    for png in PNG_DIR.glob("slide_*.png"):
        png.unlink()
    try:
        PNG_DIR.rmdir()
    except OSError:
        pass
    return True


def main() -> int:
    if not ANALYSIS.exists():
        raise SystemExit("deliverables/analysis.json not found — run "
                         "`python -m saleskpi --deliverables` first.")
    a = json.loads(ANALYSIS.read_text(encoding="utf-8"))
    n = build_pdf(a)
    print(f"wrote {OUT_PDF.relative_to(ROOT)} ({n} slides)")
    if build_pptx(a, n):
        print(f"wrote {OUT_PPTX.relative_to(ROOT)} ({n} slides)")
    else:
        print("python-pptx not installed — skipped executive_review.pptx (PDF still written)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
